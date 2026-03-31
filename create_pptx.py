from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap

# Create presentation with 16:9 aspect ratio (widescreen)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_gradient_background(slide):
    """Add a blue gradient background"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(30, 64, 175)  # Blue

def add_title_slide(prs, title, subtitle, tagline=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_gradient_background(slide)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12.333), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    if tagline:
        tag_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(10.333), Inches(0.8))
        tf = tag_box.text_frame
        p = tf.paragraphs[0]
        p.text = tagline
        p.font.size = Pt(18)
        p.font.italic = True
        p.font.color.rgb = RGBColor(200, 200, 200)
        p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide)
    
    # Title bar
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Content box
    content_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(12.333), Inches(5.6))
    content_shape.fill.solid()
    content_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    content_shape.line.fill.background()
    
    # Bullets
    text_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.8), Inches(5.2))
    tf = text_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(31, 41, 55)
        p.space_after = Pt(12)

def add_two_column_slide(prs, title, left_title, left_bullets, right_title, right_bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Left box
    left_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(6), Inches(5.6))
    left_shape.fill.solid()
    left_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    left_shape.line.fill.background()
    
    left_text = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(5.6), Inches(5.2))
    tf = left_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(31, 41, 55)
    for bullet in left_bullets:
        p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(31, 41, 55)
        p.space_after = Pt(8)
    
    # Right box
    right_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.833), Inches(1.4), Inches(6), Inches(5.6))
    right_shape.fill.solid()
    right_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    right_shape.line.fill.background()
    
    right_text = slide.shapes.add_textbox(Inches(7.033), Inches(1.6), Inches(5.6), Inches(5.2))
    tf = right_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(31, 41, 55)
    for bullet in right_bullets:
        p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(31, 41, 55)
        p.space_after = Pt(8)

# Slide 1: Title
add_title_slide(prs, "📖 StoriesOfUs", "Social Issue Storytelling Platform", 
                '"Empowering Voices, Inspiring Change: Share Your Story, Transform Communities."')

# Slide 2: Abstract
add_content_slide(prs, "ABSTRACT", [
    "StoriesOfUs is a comprehensive storytelling platform that enables users to share real-life experiences related to social issues including mental health, discrimination, inequality, and social challenges.",
    "The platform features user authentication, media uploads (images/videos), real-time chat, and a friends system for community building."
])

# Slide 3: Problem Statement
add_content_slide(prs, "PROBLEM STATEMENT", [
    "Lack of Safe Spaces: Many individuals facing social challenges lack platforms to share experiences without judgment.",
    "Awareness Gap: Social issues often remain hidden due to stigma, preventing community understanding and support.",
    "Disconnected Communities: People experiencing similar challenges often feel isolated, unaware others share their struggles.",
    "Need for Anonymity: Fear of identification prevents many from speaking up about sensitive topics."
])

# Slide 4: Key Features
add_two_column_slide(prs, "KEY FEATURES",
    "📝 Story Features", [
        "Story Submission: Share experiences with text",
        "Media Uploads: Add images & videos (up to 5 files)",
        "Category Organization: Filter by topic",
        "Anonymous Posting: Privacy protection"
    ],
    "👥 Social Features", [
        "User Authentication: Secure login/signup",
        "Friends System: Connect with other users",
        "Real-time Chat: Socket.IO powered messaging",
        "Online Status: See who's online"
    ])

# Slide 5: Categories
add_content_slide(prs, "STORY CATEGORIES", [
    "🧠 Mental Health - Stories about mental wellness and challenges",
    "✊ Discrimination - Experiences of unfair treatment",
    "⚖️ Inequality - Stories highlighting social disparities",
    "🌍 Social Challenges - General social issues and struggles",
    "📖 Other - Miscellaneous experiences worth sharing"
])

# Slide 6: User Flow
add_content_slide(prs, "USER FLOW", [
    "🏠 Visit Home → 📚 Browse Stories → 🔍 Filter by Category → 📖 Read Story → ❤️ Like",
    "✍️ Share Story → 📝 Fill Form → 🏷️ Select Category → 📁 Upload Media → ✅ Submit → 🎉 Published!",
    "🔐 Login → 👥 View Friends → 💬 Start Chat → 📤 Send Messages → 🟢 See Online Status"
])

# Slide 7: Data Model
add_two_column_slide(prs, "DATA MODEL",
    "📖 Story & Category", [
        "Story: title, content, author, category",
        "Story: featured, likes, media[] (type, url)",
        "Category: name, slug, description, icon, color"
    ],
    "👤 User & Message", [
        "User: name, email, password, friends[]",
        "Message: sender, receiver, content, read",
        "Just MONGODB for Passwords"
    ])

# Slide 8: Application Pages
add_two_column_slide(prs, "APPLICATION PAGES",
    "📄 Core Pages", [
        "🏠 Home - Featured stories & grid",
        "📚 Stories - Browse & filter",
        "📖 Story Detail - Full view + media",
        "✍️ Submit Story - Form + media upload"
    ],
    "🔐 Auth & Social Pages", [
        "🔑 Login - Email & password",
        "📝 Signup - New registration",
        "💬 Chat - Real-time messaging",
        "👥 Friends - Online status & list"
    ])

# Slide 9: Architecture
add_content_slide(prs, "SYSTEM ARCHITECTURE", [
    "🖥️ React Frontend ⟷ ⚡ Express.js API ⟷ 🔌 Socket.IO ⟷ 🗄️ MongoDB",
    "Frontend: React.js + React Router + Context API for state",
    "Backend: Node.js + Express.js + Socket.IO for real-time chat",
    "Auth: MongoDB for password storage",
    "Media: Multer for file uploads (images/videos)"
])

# Slide 10: Technology Stack
add_two_column_slide(prs, "TECHNOLOGY STACK",
    "⚛️ Frontend & Backend", [
        "React.js - Frontend + Router",
        "Node.js - Runtime",
        "Express.js - REST API",
        "MongoDB - Database"
    ],
    "🔧 Additional Tools", [
        "Socket.IO - Real-time Chat",
        "Auth - Login/Signup",
        "Multer - File Uploads",
        "Docker - Containerization"
    ])

# Slide 11: API Endpoints
add_two_column_slide(prs, "API ENDPOINTS",
    "📖 Stories & Auth", [
        "GET /api/stories - List all",
        "GET /api/stories/:id - Single story",
        "POST /api/stories - Create (+ media)",
        "POST /api/auth/register - Sign up",
        "POST /api/auth/login - Login",
        "GET /api/auth/me - Current user"
    ],
    "👥 Friends & Chat", [
        "GET /api/friends - List friends",
        "GET /api/friends/search - Search users",
        "POST /api/friends/add/:userId - Add friend",
        "GET /api/chat/:userId - Get messages",
        "POST /api/chat/:userId - Send message",
        "🔌 Socket: user_online, send_message"
    ])

# Slide 12: Implemented & Future
add_two_column_slide(prs, "IMPLEMENTED & FUTURE SCOPE",
    "✅ Implemented", [
        "User Authentication: Login/signup ✓",
        "Media Uploads: Images & videos ✓",
        "Real-time Chat: Socket.IO messaging ✓",
        "Friends System: Add & manage friends ✓",
        "Modern UI: Dark theme + glassmorphism ✓"
    ],
    "🎯 Future Enhancements", [
        "Comments: Story discussions",
        "Search: Full-text search",
        "AI Moderation: Content safety",
        "Notifications: Push alerts",
        "Mobile App: iOS & Android"
    ])

# Slide 13: References
add_content_slide(prs, "REFERENCES", [
    "React.js Documentation - https://react.dev/",
    "Node.js Documentation - https://nodejs.org/docs/",
    "Express.js Guide - https://expressjs.com/",
    "MongoDB Manual - https://www.mongodb.com/docs/manual/",
    "Mongoose ODM - https://mongoosejs.com/docs/",
    "Socket.IO Documentation - https://socket.io/docs/",
    "Multer (File Upload) - https://github.com/expressjs/multer",
    "Docker Documentation - https://docs.docker.com/"
])

# Slide 14: Conclusion
add_title_slide(prs, "Thank You! 🙏",
    "StoriesOfUs - Every story shared is a step toward a more empathetic world.",
    "Questions?")

# Save
prs.save("StoriesOfUs_Presentation.pptx")
print("✅ PPTX created: StoriesOfUs_Presentation.pptx")

